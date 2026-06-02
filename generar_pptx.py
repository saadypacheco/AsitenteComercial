# -*- coding: utf-8 -*-
"""Genera la presentación del cliente (inglés) siguiendo el modelo de apops.
Uso: python generar_pptx.py
Salida: Presentacion-Asistente.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Paleta (marca del prototipo) ----
BG     = RGBColor(0x0F, 0x14, 0x19)
PANEL  = RGBColor(0x1A, 0x21, 0x38)
PANEL2 = RGBColor(0x14, 0x18, 0x23)
A1     = RGBColor(0x66, 0x7E, 0xEA)   # indigo
A2     = RGBColor(0x76, 0x4B, 0xA2)   # purple
CYAN   = RGBColor(0x8A, 0xB4, 0xFF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTE   = RGBColor(0x9A, 0xA5, 0xB8)
DIM    = RGBColor(0x6B, 0x72, 0x80)
DANGER = RGBColor(0xFF, 0x6B, 0x6B)
WARN   = RGBColor(0xFF, 0xA9, 0x40)
OK     = RGBColor(0x52, 0xC4, 0x1A)
LINE   = RGBColor(0x2A, 0x33, 0x4D)
FONT   = "Calibri"
FONT_H = "Calibri"

prs = Presentation()
prs.slide_width  = Emu(12192000)
prs.slide_height = Emu(6858000)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]

TOTAL = 18
_n = [0]
IMG = "public/imagenes/"

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def grad(shape, c1, c2, angle=45):
    """Aplica un degradé lineal a la forma."""
    sp = shape.fill._xPr
    for tag in ('a:noFill','a:solidFill','a:gradFill','a:blipFill','a:pattFill','a:grpFill'):
        e = sp.find(qn(tag))
        if e is not None: sp.remove(e)
    g = sp.makeelement(qn('a:gradFill'), {})
    lst = g.makeelement(qn('a:gsLst'), {})
    for pos, col in ((0, c1), (100000, c2)):
        gs = g.makeelement(qn('a:gs'), {'pos': str(pos)})
        clr = g.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (col[0], col[1], col[2])})
        gs.append(clr); lst.append(gs)
    g.append(lst)
    lin = g.makeelement(qn('a:lin'), {'ang': str(int(angle*60000)), 'scaled': '1'})
    g.append(lin)
    sp.insert(0, g)

def rect(s, l, t, w, h, fill=PANEL, line=None, rounded=True, line_w=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    if rounded:
        try: shp.adjustments[0] = 0.06
        except Exception: pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, Pt(0))
    return tf

def para(tf, runs, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
         before=0, after=0, line=1.15, first=False, font=FONT):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before); p.space_after = Pt(after)
    try: p.line_spacing = line
    except Exception: pass
    if isinstance(runs, str):
        runs = [(runs, size, color, bold, font)]
    for item in runs:
        txt, *rest = item
        sz   = rest[0] if len(rest) > 0 else size
        col  = rest[1] if len(rest) > 1 else color
        bd   = rest[2] if len(rest) > 2 else bold
        fnt  = rest[3] if len(rest) > 3 else font
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bd; r.font.name = fnt
        r.font.color.rgb = col
    return p

def topbar(s):
    rect(s, 0, 0, SW, 0.10, fill=A1, rounded=False)

def footer(s, label=True):
    _n[0] += 1
    f = tb(s, 0.7, SH-0.5, SW-1.4, 0.3, anchor=MSO_ANCHOR.MIDDLE)
    para(f, [("Assistant · Proposal 2026", 9, DIM, False),
             ("        " + f"{_n[0]} / {TOTAL}", 9, DIM, False)], first=True)

def section_head(s, kicker, title, sub=None):
    topbar(s)
    tf = tb(s, 0.7, 0.55, SW-1.4, 0.4)
    para(tf, kicker.upper(), 12, CYAN, True, first=True)
    t2 = tb(s, 0.7, 0.95, SW-1.4, 0.9)
    para(t2, title, 30, WHITE, True, first=True)
    if sub:
        t3 = tb(s, 0.7, 1.75, SW-1.4, 0.5)
        para(t3, sub, 14, MUTE, False, first=True)

def bullets(s, l, t, w, items, gap=0.0, size=13, lh=1.2):
    """items: (emoji, label, desc) -> tarjeta con filas."""
    tf = tb(s, l, t, w, 5)
    for i,(emoji,label,desc) in enumerate(items):
        runs = []
        if emoji: runs.append((emoji+"  ", size+2, WHITE, False))
        if label: runs.append((label, size, WHITE, True))
        if label and desc: runs.append(("  —  ", size, DIM, False))
        if desc: runs.append((desc, size, MUTE, False))
        para(tf, runs, after=9, line=lh, first=(i==0))

def add_shot(s, path, top=1.9, max_w=11.7, max_h=3.95):
    pic = s.shapes.add_picture(path, Inches(0.5), Inches(top))
    natw, nath = pic.width, pic.height
    w = int(Inches(max_w)); h = int(w * nath / natw)
    if h > int(Inches(max_h)):
        h = int(Inches(max_h)); w = int(h * natw / nath)
    pic.width = Emu(w); pic.height = Emu(h)
    pic.left = Emu(int(prs.slide_width/2 - w/2)); pic.top = Inches(top)
    pic.line.color.rgb = LINE; pic.line.width = Pt(1.25)
    return pic

def keypoint(s, text, top=6.05):
    cb = rect(s, 0.7, top, SW-1.4, 0.6, fill=PANEL, line=A1, line_w=1.0)
    ct = cb.text_frame; ct.word_wrap=True; ct.vertical_anchor=MSO_ANCHOR.MIDDLE
    ct.margin_left=Pt(14); ct.margin_right=Pt(14)
    para(ct, [("💡 KEY POINT    ", 12, CYAN, True), (text, 12.5, WHITE, False)], first=True)

def shot_slide(kicker, title, img, kp, portrait=False):
    s = slide(); section_head(s, kicker, title)
    if portrait:
        pic = s.shapes.add_picture(IMG+img, Inches(1.6), Inches(1.95))
        natw, nath = pic.width, pic.height
        h = int(Inches(4.55)); w = int(h * natw / nath)
        pic.height = Emu(h); pic.width = Emu(w); pic.left = Inches(1.7); pic.top = Inches(2.0)
        pic.line.color.rgb = LINE; pic.line.width = Pt(1.25)
        kt = tb(s, 5.8, 3.1, 6.7, 2.5)
        para(kt, "💡 KEY POINT", 12, CYAN, True, first=True)
        para(kt, kp, 16, WHITE, False, before=10, line=1.4)
    else:
        add_shot(s, IMG+img)
        keypoint(s, kp)
    footer(s)
    return s

# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = slide()
rect(s, 0, 0, 0.18, SH, fill=A1, rounded=False)
band = rect(s, 0, 0, SW, SH, fill=PANEL2, rounded=False);
# franja de acento decorativa arriba a la derecha
deco = rect(s, SW-4.2, -1.2, 5.5, 3.0, fill=A2, rounded=True); grad(deco, A1, A2, 30)
deco.rotation = 18
tf = tb(s, 0.9, 2.0, 9.5, 0.5)
para(tf, "SALES LEADERSHIP  ·  DIGITAL ASSISTANT", 13, CYAN, True, first=True)
t = tb(s, 0.9, 2.5, 10.5, 1.3)
para(t, "Assistant", 58, WHITE, True, first=True)
tg = tb(s, 0.9, 3.85, 9.8, 1.4)
para(tg, "The operating system for sales leaders. Know what to do today — "
         "and act on your whole team from one place. Simple, from the phone.",
     19, MUTE, False, first=True, line=1.3)
rect(s, 0.95, 5.5, 2.2, 0.05, fill=A1, rounded=False)
b = tb(s, 0.9, 5.7, 10, 0.9)
para(b, "Sales-team management platform", 13, WHITE, False, first=True)
para(b, "Proposal · June 2026", 12, DIM, False, after=0)

# ============================================================
# SLIDE 2 — AGENDA
# ============================================================
s = slide(); section_head(s, "Agenda", "What we'll cover.",
                          "A walkthrough of the problem, the solution and the plan.")
agenda = [
    ("01","The problem","Why the leader's communication gets lost today"),
    ("02","The solution","What the Assistant is and how it feels"),
    ("03","See it today","The clickable prototype: leader & agent apps"),
    ("04","How it works","The loop: signals → engine → action"),
    ("05","Where we're going","Phased roadmap and what's still missing"),
    ("06","What we need","To move forward together"),
]
cols = [0.7, 6.95]; rowy = [2.35, 3.65, 4.95]
for i,(num,title,desc) in enumerate(agenda):
    cx = cols[i//3]; cy = rowy[i%3]
    nt = tb(s, cx, cy, 1.0, 1.0)
    para(nt, num, 30, A1, True, first=True)
    tt = tb(s, cx+1.0, cy, 4.9, 1.1)
    para(tt, title, 16, WHITE, True, first=True)
    para(tt, desc, 12, MUTE, False, after=0)
footer(s)

# ============================================================
# SLIDE 3 — 01 THE PROBLEM
# ============================================================
s = slide(); section_head(s, "01 · The problem", "The leader is drowning.",
                          "Managing ~1,300 agents by hand and by memory.")
bullets(s, 0.7, 2.45, SW-1.4, [
    ("📱","Too much WhatsApp","Everything comes in through one channel and important things slip through."),
    ("🚪","Who's leaving?","No early visibility on agents about to quit or who needs help first."),
    ("🐢","Slow onboarding","New agents take long to produce, and there's no view of where they get stuck."),
    ("🧠","Run from memory","Follow-ups, reminders and decisions all in the leader's head → stress, no scale."),
], size=14)
cb = rect(s, 0.7, 6.0, SW-1.4, 0.62, fill=PANEL, line=A1, line_w=1.2)
ct = cb.text_frame; ct.word_wrap=True; ct.vertical_anchor=MSO_ANCHOR.MIDDLE
ct.margin_left=Pt(14)
para(ct, [("THEY DON'T NEED MORE INFORMATION.  ", 13, CYAN, True),
          ("They need to know where to act today.", 13, WHITE, False)], first=True)
footer(s)

# ============================================================
# SLIDE 4 — BEFORE vs AFTER
# ============================================================
s = slide(); section_head(s, "01 · Before vs after", "How each flow changes.",
                          "From scattered and manual to prioritized and measured.")
# columna izquierda HOY
lc = rect(s, 0.7, 2.45, 5.7, 3.7, fill=PANEL2, line=LINE)
lt = tb(s, 0.95, 2.65, 5.2, 0.4); para(lt, "TODAY · WITHOUT THE APP", 12, DANGER, True, first=True)
bullets(s, 0.95, 3.15, 5.2, [
    ("✉️","Messages that get lost","WhatsApp floods, nothing measured"),
    ("❓","No idea who's at risk","Reacts late, loses agents"),
    ("📋","Onboarding by hand","No view of who's stuck"),
    ("🗂️","Data in the head","Forgotten follow-ups"),
], size=12.5)
rc = rect(s, 6.95, 2.45, 5.68, 3.7, fill=PANEL, line=A1, line_w=1.2)
rt = tb(s, 7.2, 2.65, 5.2, 0.4); para(rt, "WITH ASISTENTE", 12, OK, True, first=True)
bullets(s, 7.2, 3.15, 5.2, [
    ("🔔","One triaged inbox","Reply from the app, every message tracked"),
    ("🎯","A single Score per agent","See risk at a glance"),
    ("🏆","Game-like onboarding","Funnel shows where each one is"),
    ("⚡","Act in 1–2 taps","Everything auto-recorded"),
], size=12.5)
footer(s)

# ============================================================
# SLIDE 5 — 02 THE SOLUTION
# ============================================================
s = slide(); section_head(s, "02 · The solution", "What it feels like to use it.",
                          "Three ideas drive everything — plus the WhatsApp inbox at the core.")
cards = [
    ("⚡","Visible urgency","On open, see what to do now — color-prioritized 🔴🟠🟡🟢."),
    ("👆","Direct action","Reply, call, assign, recognize — like a digital wallet."),
    ("🎯","Radical simplicity","One thing at a time. Built for a non-technical user."),
]
cw = 3.9; gap = 0.18; x0 = 0.7
for i,(e,t1,d) in enumerate(cards):
    cx = x0 + i*(cw+gap)
    c = rect(s, cx, 2.4, cw, 2.0, fill=PANEL, line=LINE)
    ctf = tb(s, cx+0.25, 2.62, cw-0.5, 1.7)
    para(ctf, e, 26, WHITE, False, first=True)
    para(ctf, t1, 15, WHITE, True, before=4)
    para(ctf, d, 12, MUTE, False, before=2, line=1.25)
pl = tb(s, 0.7, 4.75, SW-1.4, 0.45)
para(pl, "THE THREE PILLARS", 12, CYAN, True, first=True)
bullets(s, 0.7, 5.2, SW-1.4, [
    ("🛡️","Control","See the whole team's status at a glance (Agent Score 0–100)."),
    ("🔭","Follow-up","Turn data into actions: who to call, who to help, who's leaving."),
    ("🚀","Development","Game-like onboarding toward the first sale."),
], size=12.5)
footer(s)

# ============================================================
# SLIDE 6 — THE AGENT SCORE
# ============================================================
s = slide(); section_head(s, "02 · The Agent Score", "One number per person.",
                          "It sums up each agent so the leader doesn't analyze 1,300 people.")
# círculo grande
circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), Inches(2.7), Inches(2.4), Inches(2.4))
circ.fill.solid(); circ.fill.fore_color.rgb = PANEL; circ.line.color.rgb = A1; circ.line.width = Pt(3)
circ.shadow.inherit=False
ctf = circ.text_frame; ctf.word_wrap=True
para(ctf, "82", 44, WHITE, True, align=PP_ALIGN.CENTER, first=True)
para(ctf, "0–100", 12, MUTE, False, align=PP_ALIGN.CENTER)
comps = [("Activity","30%"),("Training attendance","20%"),("Goal compliance","30%"),("Interaction","20%")]
ty=2.85
for name,w in comps:
    row = tb(s, 4.2, ty, 8.3, 0.5)
    para(row, [(name, 14, WHITE, True), ("    "+w, 12, A1, True)], first=True)
    bar = rect(s, 4.2, ty+0.34, 6.0, 0.10, fill=LINE, rounded=False)
    fillw = {"30%":1.8,"20%":1.2}[w]
    rect(s, 4.2, ty+0.34, fillw, 0.10, fill=A1, rounded=False)
    ty += 0.82
band = tb(s, 1.1, 5.35, 2.4, 0.4)
para(band, "🟢 HEALTHY", 13, OK, True, align=PP_ALIGN.CENTER, first=True)
footer(s)

# ============================================================
# SLIDES 7–12 — SEE IT TODAY (LIVE SCREENSHOTS)
# ============================================================
shot_slide("03 · See it today · Access", "The login.", "login.png",
           "Simple sign-in, multi-language (EN / ES / PT) and a WhatsApp-code option — on brand with how agents already communicate.")
shot_slide("03 · See it today · My Day", "What to do today, at a glance.", "dashboardCecilia.png",
           "Day timeline (what's left and when) + team summary + quick actions + a prioritized task list with a step-by-step guided mode.")
shot_slide("03 · See it today · Inbox", "WhatsApp, triaged.", "MensajesCecilia.png",
           "Every incoming message sorted into Urgent / Can wait / Answered — the leader reads and replies from one place.")
shot_slide("03 · See it today · Agents", "The whole team, by urgency.", "AgentesCecilia.png",
           "One Score (0–100) per agent, filters by status, and communication built into each card.")
shot_slide("03 · See it today · Onboarding", "Who's stuck on the way to the first sale.", "AgentesOnboardingCecilia.png",
           "A funnel of the 8 steps and a flag on every agent who got stuck — so the leader acts before they quit.")
shot_slide("03 · See it today · Agent app", "The agent's pocket coach.", "AgentesDashboard.png",
           "The agent sees their Score, today's urgent task, agenda, learning path and alerts. Game-like and mobile-first.", portrait=True)

# ============================================================
# SLIDE 9 — HOW IT WORKS (loop)
# ============================================================
s = slide(); section_head(s, "04 · How it works", "A closed loop.",
                          "Every message and action feeds the system back.")
steps = [
    ("Signals in","WhatsApp · activity · Zoom · production", A1),
    ("Engine calculates","Agent Score + prioritized alerts", A2),
    ("Leader acts","'What to do today' in 1–2 taps", A1),
    ("Recorded","History feeds back into the system", OK),
]
bw=2.85; gap=0.25; x0=0.7; y=3.1
for i,(t1,d,col) in enumerate(steps):
    cx = x0 + i*(bw+gap)
    c = rect(s, cx, y, bw, 1.9, fill=PANEL, line=col, line_w=1.4)
    ctf = tb(s, cx+0.22, y+0.25, bw-0.44, 1.5)
    para(ctf, f"{i+1}", 20, col, True, first=True)
    para(ctf, t1, 15, WHITE, True, before=2)
    para(ctf, d, 11.5, MUTE, False, before=3, line=1.2)
    if i < 3:
        ar = tb(s, cx+bw-0.02, y+0.6, 0.3, 0.6)
        para(ar, "→", 22, DIM, True, align=PP_ALIGN.CENTER, first=True)
cap = tb(s, 0.7, 5.3, SW-1.4, 0.6)
para(cap, [("The Agent Score and the alerts are computed automatically — ", 13, MUTE, False),
           ("never typed in by hand. That's the intelligence of the system.", 13, WHITE, True)],
     first=True, line=1.3)
footer(s)

# ============================================================
# SLIDE 10 — ROADMAP
# ============================================================
s = slide(); section_head(s, "05 · Where we're going", "Phased roadmap.",
                          "Value early: Phase 1 solves ~80% of the pain without AI.")
phases = [
    ("Phase 0 · Discovery","In progress","Validate data access, team structure, WhatsApp number, agenda.", WARN),
    ("Phase 1 · MVP","Next","Dashboard + Score + Inbox + Alerts + Onboarding, on a real database.", A1),
    ("Phase 2 · Automation","Later","WhatsApp API, automatic Zoom attendance, WFG data, push, mass messaging.", A1),
    ("Phase 3 · Intelligence (AI)","Future","24/7 assistant, sales simulator, churn prediction, meeting summaries.", A2),
]
y=2.45
for name,tag,desc,col in phases:
    c = rect(s, 0.7, y, SW-1.4, 0.92, fill=PANEL, line=LINE)
    nt = tb(s, 0.95, y+0.13, 4.0, 0.7)
    para(nt, name, 15, WHITE, True, first=True)
    para(nt, tag, 11, col, True, before=1)
    dt = tb(s, 5.1, y+0.13, 7.4, 0.7, anchor=MSO_ANCHOR.MIDDLE)
    para(dt, desc, 12.5, MUTE, False, first=True, line=1.2)
    rect(s, 0.7, y, 0.09, 0.92, fill=col, rounded=False)
    y += 1.02
footer(s)

# ============================================================
# SLIDE 11 — WHAT'S MISSING
# ============================================================
s = slide(); section_head(s, "05 · What's still missing", "Honest about the road ahead.",
                          "Presented as next phases, not gaps.")
bullets(s, 0.7, 2.45, SW-1.4, [
    ("🔌","Real integrations","WhatsApp Business API, Zoom attendance, WFG production data, push."),
    ("⚙️","The backend","Database, the engine that computes Score & alerts, and security (US data region)."),
    ("📱","Platform decision","App stores vs installable web app — defines the technical path."),
    ("🌐","Full multi-language","Login already shows it; extend across the whole app (English by default)."),
], size=14)
footer(s)

# ============================================================
# SLIDE 12 — WHAT WE NEED FROM YOU
# ============================================================
s = slide(); section_head(s, "06 · What we need from you", "To close Phase 0.",
                          "Five questions that unlock the whole design.")
qs = [
    ("1","What does your day look like, and how do you manage your agenda today?"),
    ("2","How are your ~1,300 agents organized? (by leader, cohort, product…)"),
    ("3","Which WhatsApp number do agents write to today?"),
    ("4","What data does WFG give you, and how do you access it?"),
    ("5","How do you onboard a new agent, and what data do you have on each?"),
]
y=2.5
for num,q in qs:
    nt = tb(s, 0.8, y, 0.7, 0.6)
    para(nt, num, 22, A1, True, first=True)
    qt = tb(s, 1.55, y+0.04, SW-2.3, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    para(qt, q, 15, WHITE, False, first=True, line=1.2)
    y += 0.78
footer(s)

# ============================================================
# SLIDE 13 — MANIFESTO
# ============================================================
s = slide()
rect(s, 0, 0, 0.18, SH, fill=A1, rounded=False)
q = tb(s, 1.2, 2.2, 10.9, 2.2)
para(q, "This isn't a CRM or a course platform.", 30, WHITE, True, first=True, line=1.15)
para(q, "It's the operating system a sales leader needs to stop drowning — "
        "and to scale a team without losing control.", 20, MUTE, False, before=10, line=1.3)
ln = tb(s, 1.2, 4.9, 10.9, 0.8)
para(ln, "We start with what hurts most and delivers value fast. The rest is an expansion, with a clear path.",
     15, CYAN, True, first=True, line=1.3)
footer(s)

# ============================================================
# SLIDE 14 — THANK YOU
# ============================================================
s = slide()
deco = rect(s, SW-4.6, SH-3.2, 6.0, 3.2, fill=A2, rounded=True); grad(deco, A1, A2, 30); deco.rotation=12
t = tb(s, 0.9, 2.4, 10, 1.2)
para(t, "Thank you.", 48, WHITE, True, first=True)
st = tb(s, 0.9, 3.6, 10, 0.5)
para(st, "Questions, comments, next steps.", 18, MUTE, False, first=True)
rect(s, 0.95, 4.5, 2.2, 0.05, fill=A1, rounded=False)
c = tb(s, 0.9, 4.75, 10, 1.5)
para(c, [("Contact   ", 13, DIM, False), ("Saady Pacheco", 15, WHITE, True)], first=True)
para(c, "saadypacheco@gmail.com", 13, CYAN, False, before=4)
para(c, "Clickable prototype: login.html · dashboard-v5-mi-dia.html · agente.html", 12, MUTE, False, before=8)
footer(s)

import sys
out = sys.argv[1] if len(sys.argv) > 1 else "Presentacion-Asistente.pptx"
prs.save(out)
print("OK ->", out, "·", len(prs.slides._sldIdLst), "slides")
